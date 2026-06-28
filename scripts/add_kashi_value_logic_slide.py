from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
SOURCE = Path(
    "/Users/bruce/Library/Containers/com.tencent.xinWeChat/Data/Documents/"
    "xwechat_files/qq738872843_7fbb/temp/drag/寻境喀什.pptx"
)
OUT = ROOT / "docs/01_产品规划/寻境喀什_补充项目逻辑页.pptx"
PREVIEW = ROOT / "tmp/寻境喀什_项目逻辑页_preview.png"
BG = ROOT / "tmp/寻境喀什_项目逻辑页_bg.png"
VISUAL = ROOT / "assets/references/web_喀什/ChatGPT Image 2026年6月21日 00_59_18 (1).png"
POLICY_SCREENSHOT = ROOT / "mct_ai_tourism_policy_2026.png"
POLICY_RD_SCREENSHOT = ROOT / "mct_ai_tourism_rd_2026.png"
POLICY_RESEARCH_SCREENSHOT = ROOT / "mct_ai_tourism_research_2026.png"
WEB_SCREENSHOT = ROOT / "assets/references/web_喀什/ChatGPT Image 2026年6月21日 00_59_18 (1).png"
APP_HOME = ROOT / "assets/references/APP/参考图/首页.png"
APP_AI = ROOT / "assets/references/APP/参考图/扫码后AI.png"
APP_TRAVEL = ROOT / "assets/references/APP/参考图/跟着游记游喀什.png"


EMU_PER_INCH = 914400

COLORS = {
    "bg": "F7F3E8",
    "brown": "122235",
    "body": "5E6870",
    "muted": "918A7C",
    "gold": "C9793E",
    "teal": "2F8F8A",
    "deep_teal": "0D6F6B",
    "card": "FFF8E9",
    "line": "D3B16E",
    "white": "FFFFFF",
}


def rgb(hex_value: str) -> RGBColor:
    hex_value = hex_value.strip("#")
    return RGBColor(
        int(hex_value[0:2], 16),
        int(hex_value[2:4], 16),
        int(hex_value[4:6], 16),
    )


def emu_to_px(value: int, total_emu: int, total_px: int) -> int:
    return round(value / total_emu * total_px)


def make_background() -> None:
    """Create a muted Kashgar visual background that does not compete with text."""
    PREVIEW.parent.mkdir(parents=True, exist_ok=True)
    BG.parent.mkdir(parents=True, exist_ok=True)

    width, height = 1672, 941
    bg = Image.new("RGB", (width, height), "#F7F3E8")

    if VISUAL.exists():
        src = Image.open(VISUAL).convert("RGB")
        # Keep the Kashgar city and mascot area, but avoid the large left headline.
        crop = src.crop((700, 90, src.width, 650))
        crop = crop.resize((720, 455), Image.LANCZOS)
        layer = Image.new("RGBA", (width, height), (0, 0, 0, 0))
        layer.paste(crop.convert("RGBA"), (900, 88))

        # Fade the visual toward the left and bottom so PPT text remains readable.
        mask = Image.new("L", (720, 455), 0)
        md = ImageDraw.Draw(mask)
        for x in range(720):
            alpha = int(55 * (x / 720) ** 0.7)
            md.line([(x, 0), (x, 455)], fill=alpha)
        for y in range(455):
            fade = int(255 * min(1, (455 - y) / 180))
            if fade < 255:
                for x in range(720):
                    mask.putpixel((x, y), int(mask.getpixel((x, y)) * fade / 255))
        layer.putalpha(Image.eval(layer.getchannel("A"), lambda a: a))
        layer_alpha = Image.new("L", (width, height), 0)
        layer_alpha.paste(mask, (900, 88))
        layer.putalpha(layer_alpha)
        bg = Image.alpha_composite(bg.convert("RGBA"), layer).convert("RGB")

    draw = ImageDraw.Draw(bg, "RGBA")
    # Soft ornamental bands echo the source deck without copying a full slide.
    draw.rectangle((0, 0, width, 74), fill=(247, 243, 232, 235))
    draw.line((70, 80, width - 70, 80), fill=(201, 121, 62, 90), width=2)
    draw.line((72, height - 68, width - 72, height - 68), fill=(145, 138, 124, 95), width=2)
    for i in range(7):
        x = -80 + i * 52
        draw.arc((x, height - 155, x + 210, height + 54), 198, 350, fill=(201, 121, 62, 55), width=3)
    for i in range(6):
        x = width - 360 + i * 48
        draw.arc((x, 28, x + 180, 178), 15, 180, fill=(47, 143, 138, 40), width=2)
    bg.save(BG)


def set_text(shape, text, font_size, color, bold=False, font="Arial", align=None):
    shape.text_frame.clear()
    shape.text_frame.word_wrap = True
    shape.text_frame.margin_left = Inches(0.02)
    shape.text_frame.margin_right = Inches(0.02)
    shape.text_frame.margin_top = Inches(0.01)
    shape.text_frame.margin_bottom = Inches(0.01)
    p = shape.text_frame.paragraphs[0]
    if align:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)


def add_text(slide, x, y, w, h, text, size, color="body", bold=False, font="Arial", align=None):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    set_text(shape, text, size, COLORS[color], bold, font, align)
    return shape


def add_card(slide, x, y, w, h, fill="card", line="line", radius=True):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(COLORS[fill])
    shape.line.color.rgb = rgb(COLORS[line])
    shape.line.width = Pt(1.1)
    return shape


def add_image_contain(slide, path: Path, x, y, w, h):
    with Image.open(path) as image:
        image_w, image_h = image.size
    box_ratio = w / h
    image_ratio = image_w / image_h
    if image_ratio >= box_ratio:
        pic_w = w
        pic_h = w / image_ratio
        pic_x = x
        pic_y = y + (h - pic_h) / 2
    else:
        pic_h = h
        pic_w = h * image_ratio
        pic_x = x + (w - pic_w) / 2
        pic_y = y
    return slide.shapes.add_picture(str(path), Inches(pic_x), Inches(pic_y), width=Inches(pic_w), height=Inches(pic_h))


def add_pill(slide, x, y, w, h, text, fill="teal", color="white", size=10.5):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(COLORS[fill])
    shape.line.color.rgb = rgb(COLORS["gold"])
    shape.line.width = Pt(0.7)
    shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_text(shape, text, size, COLORS[color], True, "Arial", PP_ALIGN.CENTER)
    return shape


def add_bullet_list(slide, x, y, w, lines, color="body", size=9.8, leading=1.0):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.42 * len(lines) + 0.1))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.03)
    tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.0)
    tf.margin_bottom = Inches(0.0)
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Arial"
        p.font.size = Pt(size)
        p.font.color.rgb = rgb(COLORS[color])
        p.space_after = Pt(leading)
    return box


def add_customer(slide, x, y, title, body, accent="teal"):
    add_pill(slide, x, y, 1.2, 0.28, title, fill=accent, size=8.5)
    add_text(slide, x + 1.34, y - 0.01, 2.25, 0.46, body, 8.7, "body", False)


def add_arrow(slide, x1, y1, x2, y2):
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    connector.line.color.rgb = rgb(COLORS["gold"])
    connector.line.width = Pt(1.4)
    return connector


def add_footer(slide, page_num="03"):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(0.61),
        Inches(7.0),
        Inches(12.58),
        Inches(7.0),
    )
    line.line.color.rgb = rgb(COLORS["muted"])
    line.line.width = Pt(0.65)
    add_text(slide, 0.61, 7.05, 2.1, 0.2, "项目逻辑总览", 9, "muted")
    add_text(slide, 12.48, 7.05, 0.38, 0.2, page_num, 9, "muted", align=PP_ALIGN.RIGHT)


def walk_shapes(shapes):
    for shape in shapes:
        yield shape
        if getattr(shape, "shape_type", None) == 6:
            yield from walk_shapes(shape.shapes)


def replace_text(slide, replacements):
    for shape in walk_shapes(slide.shapes):
        if not hasattr(shape, "text") or not shape.text.strip():
            continue
        source = " ".join(shape.text.split())
        if source not in replacements:
            continue
        target = replacements[source]
        tf = shape.text_frame
        if not tf.paragraphs:
            shape.text = target
            continue
        paragraph = tf.paragraphs[0]
        if paragraph.runs:
            paragraph.runs[0].text = target
            for run in paragraph.runs[1:]:
                run.text = ""
        else:
            paragraph.text = target
        for extra in tf.paragraphs[1:]:
            extra.text = ""


def replace_text_in_deck(prs: Presentation, replacements):
    for slide in prs.slides:
        replace_text(slide, replacements)


def rewrite_opportunity_slide(prs: Presentation):
    slide = prs.slides[1]
    replace_text(
        slide,
        {
            "地方文旅资源很丰富，下一步需要更好的运营连接。": "为什么现在可以把文旅做成长期资产？",
            "地方有历史、有非遗、有景区、有图书、有影像，下一步可以把它们更稳定地连接到游客服务、内容传播、运营数据和商业线索。": "AI、二维码、小程序和内容后台已经成熟，一次性项目可以变成持续服务和数据资产。",
            "资源协同": "AI 可用",
            "资料、图片、视频、图书散落，无法被统一检索和调用。": "官方资料能转成可问、可读、可生成的服务。",
            "服务连续性": "入口轻",
            "导游、讲解、问答靠人力，游客离场后服务就结束。": "小程序、二维码、展陈码、图书码适合进入真实场景。",
            "内容复用": "预算清",
            "宣传内容靠活动和外包，难以持续生产和复用。": "数字文旅、智慧景区、公共文化、宣传运营都有交付口径。",
            "数据沉淀": "可复制",
            "扫码、提问、路线、作品没有沉淀，运营决策仍然靠感觉。": "喀什跑通后，可复制到景区、馆、图书、研学和公益活动。",
            "改进空间": "机会判断",
            "让游客体验更连贯，让机构资产更清晰，让运营成效更容易复盘。": "现在适合做的是一套内容资产运营系统。",
            "这也是内容资产运营系统适合进入文旅项目预算的原因。": "下一页说明：为什么资源过去没留住，以及星河寻境怎么沉淀。",
            "现状与机会": "机会窗口",
        },
    )


def add_opportunity_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    slide.shapes.add_picture(str(BG), 0, 0, width=prs.slide_width, height=prs.slide_height)

    add_text(slide, 0.58, 0.34, 1.4, 0.2, "WHY NOW", 9.5, "gold", True)
    add_text(
        slide,
        0.58,
        0.62,
        8.6,
        0.46,
        "为什么现在可以把文旅做成长期资产？",
        27,
        "brown",
        True,
        "Georgia",
    )
    add_text(
        slide,
        0.6,
        1.16,
        8.8,
        0.3,
        "AI 能理解内容，二维码和小程序能进入现场，后台能持续运营，文旅项目正在从“做完交付”转向“长期资产”。",
        11.4,
        "body",
    )

    cards = [
        (
            "01",
            "AI 可用",
            "文本、视觉、语音、检索能力成熟，官方资料可以转成可问、可读、可生成的服务。",
            "teal",
        ),
        (
            "02",
            "入口够轻",
            "小程序、二维码、展陈码、图书码已经是低成本入口，适合进入真实文旅场景。",
            "deep_teal",
        ),
        (
            "03",
            "采购有口径",
            "数字文旅、智慧景区、公共文化、宣传运营，都能对应清晰交付口径。",
            "gold",
        ),
        (
            "04",
            "复制空间大",
            "一个地方跑通后，可复制到景区、博物馆、图书、研学和公益活动。",
            "gold",
        ),
    ]
    x_positions = [0.72, 3.82, 6.92, 10.02]
    tag_labels = ["技术基础", "现场入口", "采购路径", "复制模型"]
    for x, (num, title, body, accent), tag in zip(x_positions, cards, tag_labels):
        add_card(slide, x, 2.03, 2.55, 3.08)
        add_text(slide, x + 0.24, 2.31, 0.44, 0.34, num, 15, accent, True, "Georgia")
        add_text(slide, x + 0.72, 2.34, 1.5, 0.28, title, 16, "brown", True, "Arial")
        add_pill(slide, x + 0.3, 2.92, 1.72, 0.3, tag, fill=accent, size=8.3)
        add_text(slide, x + 0.32, 3.38, 1.92, 1.12, body, 10.2, "body")

    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.72), Inches(5.77), Inches(11.9), Inches(0.72))
    bar.fill.solid()
    bar.fill.fore_color.rgb = rgb(COLORS["brown"])
    bar.line.color.rgb = rgb(COLORS["gold"])
    bar.line.width = Pt(1)
    bar.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_text(
        bar,
        "结论：现在适合做的不是一个 AI 演示，而是一套能持续沉淀内容、服务和数据的运营系统。",
        14.0,
        COLORS["white"],
        True,
        "Arial",
        PP_ALIGN.CENTER,
    )

    add_footer(slide, "02")
    return slide


def add_policy_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    slide.shapes.add_picture(str(BG), 0, 0, width=prs.slide_width, height=prs.slide_height)

    add_text(slide, 0.58, 0.34, 1.65, 0.2, "POLICY SIGNAL", 9.5, "gold", True)
    add_text(
        slide,
        0.58,
        0.62,
        8.6,
        0.46,
        "2026 年，“人工智能+文化和旅游”进入应用试点",
        25,
        "brown",
        True,
        "Georgia",
    )
    add_text(
        slide,
        0.6,
        1.16,
        8.3,
        0.3,
        "文旅部已把 AI 文旅从方向倡议推进到试点申报、科技创新和数智化研究，项目需要对应“数据治理、场景应用、标准与成效”。",
        10.8,
        "body",
    )

    policy_items = [
        (
            POLICY_SCREENSHOT,
            "应用试点",
            "“人工智能+文化和旅游”应用试点",
            "2026-03-25",
        ),
        (
            POLICY_RD_SCREENSHOT,
            "科技创新",
            "2026 年国家文旅科技创新项目",
            "2026-03-24",
        ),
        (
            POLICY_RESEARCH_SCREENSHOT,
            "数智研究",
            "2026 年文旅数智化研究课题",
            "2026-04-13",
        ),
    ]
    xs = [0.74, 4.72, 8.7]
    for x, (path, tag, title, date) in zip(xs, policy_items):
        add_card(slide, x, 1.72, 3.55, 3.95, fill="white", line="line")
        if path.exists():
            add_image_contain(slide, path, x + 0.16, 1.9, 3.22, 2.24)
        add_pill(slide, x + 0.24, 4.35, 1.05, 0.3, tag, "deep_teal" if tag != "数智研究" else "gold", size=8.1)
        add_text(slide, x + 1.42, 4.35, 1.75, 0.24, date, 8.4, "muted", align=PP_ALIGN.RIGHT)
        add_text(slide, x + 0.24, 4.78, 2.86, 0.42, title, 10.0, "brown", True)
        add_text(slide, x + 0.24, 5.26, 2.9, 0.22, "来源：中华人民共和国文化和旅游部官网", 7.6, "muted")

    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(6.08), Inches(12.05), Inches(0.58))
    bar.fill.solid()
    bar.fill.fore_color.rgb = rgb(COLORS["brown"])
    bar.line.color.rgb = rgb(COLORS["gold"])
    bar.line.width = Pt(1)
    bar.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_text(
        bar,
        "判断：政策侧已经从方向倡议进入试点、创新和研究部署；星河寻境要对齐的是场景应用、数据治理和可评估成果。",
        12.6,
        COLORS["white"],
        True,
        "Arial",
        PP_ALIGN.CENTER,
    )

    add_footer(slide, "02")
    return slide


def add_value_logic_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    slide.shapes.add_picture(str(BG), 0, 0, width=prs.slide_width, height=prs.slide_height)

    add_text(slide, 0.58, 0.34, 1.7, 0.2, "VALUE LOGIC", 9.5, "gold", True)
    add_text(
        slide,
        0.58,
        0.62,
        9.0,
        0.46,
        "地方文旅的关键问题：有资源，但没有沉淀成可运营资产",
        24.5,
        "brown",
        True,
        "Georgia",
    )
    add_text(
        slide,
        0.6,
        1.14,
        8.6,
        0.28,
        "景区、人文、图书、影像和活动都存在，但常停留在一次体验、一次传播、一次项目里，难以被地方持续调用。",
        11.2,
        "body",
    )

    # Left resource card.
    add_card(slide, 0.62, 1.7, 3.2, 3.86)
    add_pill(slide, 0.91, 1.98, 1.58, 0.33, "问题本质", "deep_teal")
    add_text(slide, 0.92, 2.43, 2.65, 0.45, "不是资源少，而是资产留不下。", 11.0, "brown", True)
    add_bullet_list(
        slide,
        0.92,
        2.96,
        2.54,
        [
            "资料、图片、视频、图书散在不同项目里，难统一调用。",
            "游客到访后互动结束，用户内容和反馈很难回到地方。",
            "宣传热度多留在外部平台，地方缺少长期内容库。",
            "活动、研学、出版做完后，素材和数据难复用。",
        ],
        size=8.65,
    )
    add_text(
        slide,
        0.92,
        4.9,
        2.48,
        0.34,
        "结果是每次都在重新组织内容、重新找素材、重新做传播。",
        9.2,
        "muted",
    )

    # Center platform card.
    add_card(slide, 4.34, 1.52, 3.95, 4.21, fill="deep_teal", line="gold")
    add_text(slide, 4.69, 1.87, 3.22, 0.35, "星河寻境", 21, "white", True, "Georgia", PP_ALIGN.CENTER)
    add_text(slide, 4.65, 2.25, 3.3, 0.25, "地方文旅 AI 内容资产平台", 10.8, "gold", True, "Arial", PP_ALIGN.CENTER)
    add_card(slide, 4.72, 2.82, 3.18, 0.58, fill="card", line="line")
    add_text(slide, 4.91, 2.96, 2.78, 0.25, "资产底座：知识库 + 影像库 + 作品库", 9.3, "brown", True, align=PP_ALIGN.CENTER)
    add_card(slide, 4.72, 3.58, 3.18, 0.58, fill="card", line="line")
    add_text(slide, 4.91, 3.72, 2.78, 0.25, "场景入口：AI 旅伴 + 景点码 + 图书码", 9.3, "brown", True, align=PP_ALIGN.CENTER)
    add_card(slide, 4.72, 4.34, 3.18, 0.58, fill="card", line="line")
    add_text(slide, 4.91, 4.48, 2.78, 0.25, "运营沉淀：问答、路线、作品、报告", 9.3, "brown", True, align=PP_ALIGN.CENTER)
    add_text(
        slide,
        4.74,
        5.14,
        3.15,
        0.22,
        "目标是把每一次游览、阅读和活动，沉淀回地方资产。",
        8.8,
        "white",
        False,
        align=PP_ALIGN.CENTER,
    )

    # Right buyer card.
    add_card(slide, 8.83, 1.7, 3.88, 3.86)
    add_pill(slide, 9.12, 1.98, 1.82, 0.33, "为什么买单", "deep_teal")
    add_text(slide, 9.13, 2.43, 2.95, 0.45, "买的不是 AI 功能，而是自己能带走、能评估、能继续用的项目成果。", 9.5, "brown", True)
    add_customer(slide, 9.08, 2.98, "文旅局", "把城市资源做成数字资产和示范成果。")
    add_customer(slide, 9.08, 3.53, "景区 / 馆", "把讲解服务、游客内容和运营数据留在自己手里。")
    add_customer(slide, 9.08, 4.08, "出版社 / 中图", "把图书和地图变成可互动、可统计的数字内容。", accent="gold")
    add_customer(slide, 9.08, 4.63, "公益 / 研学", "把活动过程、学习成果和报告标准化复用。", accent="gold")

    add_arrow(slide, 3.93, 3.65, 4.25, 3.65)
    add_arrow(slide, 8.41, 3.65, 8.73, 3.65)

    # Bottom buying reason.
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.62), Inches(6.04), Inches(12.08), Inches(0.66))
    bar.fill.solid()
    bar.fill.fore_color.rgb = rgb(COLORS["brown"])
    bar.line.color.rgb = rgb(COLORS["gold"])
    bar.line.width = Pt(1)
    bar.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    set_text(
        bar,
        "一句话：把一次性资源、一次性流量、一次性活动，变成地方长期拥有的内容资产和运营资产。",
        13.4,
        COLORS["white"],
        True,
        "Arial",
        PP_ALIGN.CENTER,
    )

    add_footer(slide)
    return slide


def add_web_example_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    slide.shapes.add_picture(str(BG), 0, 0, width=prs.slide_width, height=prs.slide_height)

    add_text(slide, 0.58, 0.34, 1.95, 0.2, "PRODUCT SAMPLE", 9.5, "gold", True)
    add_text(slide, 0.58, 0.62, 7.6, 0.46, "喀什网页示例：产品入口与能力呈现", 25, "brown", True, "Georgia")
    add_text(slide, 0.6, 1.16, 7.9, 0.3, "网页端承担对外展示、项目沟通和产品能力说明，把地方文旅资产平台讲成可看、可点、可复制的产品入口。", 10.8, "body")

    add_card(slide, 0.62, 1.62, 8.05, 4.75, fill="white", line="line")
    if WEB_SCREENSHOT.exists():
        add_image_contain(slide, WEB_SCREENSHOT, 0.75, 1.76, 7.78, 4.45)

    add_card(slide, 9.03, 1.68, 3.55, 4.62)
    add_pill(slide, 9.34, 1.98, 1.72, 0.33, "页面价值", "deep_teal")
    add_text(slide, 9.34, 2.46, 2.48, 0.32, "面向客户的产品入口呈现。", 11.0, "brown", True)
    add_bullet_list(
        slide,
        9.34,
        2.96,
        2.55,
        [
            "呈现喀什首页、能力模块和三大入口。",
            "将知识库、影像库、AI 旅伴、扫码伴读和内容运营前台化。",
            "用于项目沟通、方案展示和交付范围对齐。",
        ],
        size=9.2,
    )
    add_text(slide, 9.34, 5.38, 2.55, 0.28, "定位：Web 端解决“看得懂、愿意试、便于汇报”。", 9.1, "muted")

    add_footer(slide, "13")
    return slide


def add_mini_program_example_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    slide.shapes.add_picture(str(BG), 0, 0, width=prs.slide_width, height=prs.slide_height)

    add_text(slide, 0.58, 0.34, 1.95, 0.2, "PRODUCT SAMPLE", 9.5, "gold", True)
    add_text(slide, 0.58, 0.62, 7.8, 0.46, "喀什小程序示例：现场服务与内容沉淀入口", 25, "brown", True, "Georgia")
    add_text(slide, 0.6, 1.16, 8.0, 0.3, "小程序端面向游客现场使用，把“扫一扫、AI 旅伴、跟着游记、记录旅行”组织成连续的服务路径。", 10.8, "body")

    phone_specs = [
        (APP_HOME, 0.78, "首页入口"),
        (APP_AI, 3.35, "AI 旅伴"),
        (APP_TRAVEL, 5.92, "跟着游记"),
    ]
    for path, x, label in phone_specs:
        add_card(slide, x, 1.72, 2.08, 4.74, fill="white", line="line")
        if path.exists():
            add_image_contain(slide, path, x + 0.1, 1.86, 1.88, 4.18)
        add_text(slide, x + 0.25, 6.12, 1.58, 0.22, label, 9.2, "brown", True, align=PP_ALIGN.CENTER)

    add_card(slide, 8.78, 1.72, 3.82, 4.74)
    add_pill(slide, 9.1, 2.02, 1.72, 0.33, "产品路径", "deep_teal")
    add_text(slide, 9.1, 2.5, 2.85, 0.32, "小程序承接游客现场使用和内容回流。", 10.8, "brown", True)
    add_bullet_list(
        slide,
        9.1,
        2.98,
        2.78,
        [
            "首页聚合三入口：扫码、跟着游记、记录旅行。",
            "AI 旅伴承接讲解、问答、路线和打卡地推荐。",
            "游记内容可以反哺路线、作品库和传播素材。",
            "后台对应二维码、知识库、素材库和用户作品管理。",
        ],
        size=9.0,
    )
    add_text(slide, 9.1, 5.56, 2.8, 0.28, "定位：小程序端承接现场服务、内容生成和数据回流。", 9.0, "muted")

    add_footer(slide, "14")
    return slide


def move_last_slide_to(prs: Presentation, index: int) -> None:
    sld_id_lst = prs.slides._sldIdLst
    last = sld_id_lst[-1]
    sld_id_lst.remove(last)
    sld_id_lst.insert(index, last)


def delete_slide(prs: Presentation, index: int) -> None:
    sld_id_lst = prs.slides._sldIdLst
    slide_id = sld_id_lst[index]
    rel_id = slide_id.rId
    prs.part.drop_rel(rel_id)
    sld_id_lst.remove(slide_id)


def polish_external_wording(prs: Presentation) -> None:
    replace_text_in_deck(
        prs,
        {
            "MVP 只预留商家字段，不做在线支付、团购、商家入驻和佣金结算，避免早期产品重心跑偏。": "首期预留本地服务字段，暂不展开在线支付、团购、商家入驻和佣金结算，聚焦内容资产与游客服务闭环。",
            "MVP 0-2 个月：喀什体验验证版": "阶段一 0-2 个月：喀什体验启动版",
            "M1 第 3-4 个月：可试点版": "阶段二 第 3-4 个月：项目试运行版",
            "支持一个地方或景区试点，具备内容录入、二维码绑定和体验反馈。": "支持一个地方或景区先行落地，具备内容录入、二维码绑定和体验反馈。",
            "M2 第 5-8 个月：地方内容运营版": "阶段三 第 5-8 个月：地方内容运营版",
            "M3 第 9-12 个月：多地复制与轻量商家服务": "阶段四 第 9-12 个月：多地复制与轻量商家服务",
            "M4 第 13-18 个月：运营平台版": "阶段五 第 13-18 个月：运营平台版",
            "M5 第 19-24 个月：平台规模化版": "阶段六 第 19-24 个月：平台规模化版",
            "阶段目标按交付成熟度递进，先证明样板，再扩展运营和规模化能力。": "阶段目标按交付成熟度递进，先跑通喀什样板，再扩展运营和规模化能力。",
            "试点验收指标：客户买完以后，要看得见结果": "项目成效指标：客户买完以后，要看得见结果",
            "证明资产底座已经建成": "呈现资产底座建设成果",
            "证明服务进入真实场景": "呈现场景入口覆盖情况",
            "证明游客真的在使用": "呈现游客真实使用情况",
            "证明内容生产效率提升": "呈现内容生产效率提升",
            "验收逻辑：先证明资产建成，再证明有人使用，最后证明可以持续运营和复制。": "成效逻辑：先看资产建成，再看真实使用，最后看持续运营和复制能力。",
            "这套指标让客户内部汇报有依据，也让第二期扩容有预算理由。": "这套指标让项目汇报有依据，也为第二期扩容形成预算理由。",
            "试点验收指标": "项目成效指标",
            "客户买单的关键不是“装系统”，而是 90 天内看到资产、入口、体验和数据闭环。": "客户买单的关键不是“装系统”，而是 90 天内看到资产、入口、体验和数据闭环。",
            "验收复盘": "成效复盘",
            "输出：验收报告": "输出：成效报告",
            "交付物越明确，客户越容易采购：资产包、场景码、AI 体验、运营数据、验收报告缺一不可。": "交付物越明确，客户越容易采购：资产包、场景码、AI 体验、运营数据、成效报告缺一不可。",
            "为什么客户会买单：买的是“可验收的文旅资产项目”": "为什么客户会买单：买的是“可交付的文旅资产项目”",
            "采购不是为 AI 试新": "采购不是为 AI 尝鲜",
            "有验收物": "有交付物",
            "一句话：把一次性宣传预算，升级成可验收、可复用、可运营的地方文旅数字资产。": "一句话：把一次性宣传预算，升级成可交付、可复用、可运营的地方文旅数字资产。",
            "首期可验收资产包和样板工程价值": "首期可交付资产包和样板工程价值",
            "把客户愿意买单的价值拆成可采购、可验收、可续费的服务包。": "把客户愿意买单的价值拆成可采购、可交付、可续费的服务包。",
            "M3 以后再逐步接入本地吃喝玩乐购服务。": "后续阶段再逐步接入本地吃喝玩乐购服务。",
        },
    )


def update_footer_numbers(prs: Presentation) -> None:
    def walk(shapes, slide_index):
        for shape in shapes:
            if hasattr(shape, "text") and shape.text.strip().isdigit():
                if shape.left > 11_000_000 and shape.top > 6_000_000:
                    shape.text = f"{slide_index:02d}"
                    shape.left = Inches(12.48)
                    shape.width = Inches(0.38)
            if getattr(shape, "shape_type", None) == 6:
                walk(shape.shapes, slide_index)

    for idx, slide in enumerate(prs.slides, start=1):
        walk(slide.shapes, idx)


def render_preview() -> None:
    """Create a close visual preview of the inserted slide."""
    canvas = Image.open(BG).convert("RGB")
    draw = ImageDraw.Draw(canvas, "RGBA")
    w, h = canvas.size

    def px(x, y):
        return round(x / 13.333333 * w), round(y / 7.5 * h)

    def rect(x, y, ww, hh, fill, outline="#D3B16E", radius=12):
        x1, y1 = px(x, y)
        x2, y2 = px(x + ww, y + hh)
        draw.rounded_rectangle((x1, y1, x2, y2), radius=radius, fill=fill, outline=outline, width=2)

    # Simple geometry preview; detailed text is verified in the PPT itself.
    rect(0.62, 1.7, 3.2, 3.86, "#FFF8E9")
    rect(4.34, 1.52, 3.95, 4.21, "#0D6F6B")
    rect(8.83, 1.7, 3.88, 3.86, "#FFF8E9")
    rect(0.62, 6.04, 12.08, 0.66, "#122235", "#C9793E", 16)
    canvas.save(PREVIEW)


def main() -> None:
    make_background()
    prs = Presentation(SOURCE)
    add_policy_slide(prs)
    move_last_slide_to(prs, 1)
    add_opportunity_slide(prs)
    move_last_slide_to(prs, 2)
    add_value_logic_slide(prs)
    move_last_slide_to(prs, 3)
    add_web_example_slide(prs)
    move_last_slide_to(prs, 13)
    add_mini_program_example_slide(prs)
    move_last_slide_to(prs, 14)
    delete_slide(prs, 4)
    polish_external_wording(prs)
    update_footer_numbers(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    render_preview()
    print(OUT)
    print(PREVIEW)


if __name__ == "__main__":
    main()
